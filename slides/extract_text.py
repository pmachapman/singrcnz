from pptx import Presentation
import glob

def clean(text):
    text = text.replace('   ', '\t')
    text = text.replace(' \t', ' ')
    text = text.replace('  ', ' ')
    text = text.replace('\r\n\r\n', '\r\n')
    return text

f = open('extract.txt', 'w')
for eachfile in glob.glob('*.pptx'):
    prs = Presentation(eachfile)
    print('----------------------', file=f)
    print(eachfile, file=f)
    print('----------------------', file=f)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if shape.text[:1] == '[' or shape.text[:4] == 'Sing':
                    continue
                else:
                    print(clean(shape.text), file=f)